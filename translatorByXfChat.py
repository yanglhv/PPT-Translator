from sparkdesk_api.core import SparkAPI
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

LANGUAGE_CODE_TO_LANGUAGE_ID = {
    '中文': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    '英文': MSO_LANGUAGE_ID.ENGLISH_US,
    '日语': MSO_LANGUAGE_ID.JAPANESE
}


class TranslatorByXfChat:
    def __init__(self, app_id, api_key, api_secret, version, target_language):
        self.app_id = app_id
        self.api_key = api_key
        self.api_secret = api_secret
        self.version = version
        self.target_language = target_language
        self.sparkAPI = SparkAPI(app_id=app_id, api_secret=api_secret, api_key=api_key, version=version)

    def translate_text(self, text):
        # sparkAPI = SparkAPI(app_id=self.app_id, api_secret=self.api_secret, api_key=self.api_key, version=self.version)
        lg_to = self.target_language
        promotTemplate = f"将输入内容翻译成{lg_to}；如果输入的是{lg_to}或者遇到了你不会翻译的词汇组合，则直接输出原输入的内容，不要处理标点符号。仅输出翻译的结果，除了翻译结果，不要输出任何其他话。现在开始："
        # print(promotTemplate)
        # history = [{
        #     "role": "system",
        #     "content": promotTemplate
        # }]
        # translated_text = self.sparkAPI.chat(query=text, history=history, temperature=0.1)
        translated_text = self.sparkAPI.chat(query=promotTemplate+text, temperature=0.4)
        return translated_text

    def translate_presentation(self, presentation: Presentation):
        slide_number = 1
        for slide in presentation.slides:
            print('Slide {slide_number} of {number_of_slides}'.format(
                slide_number=slide_number,
                number_of_slides=len(presentation.slides)))
            slide_number += 1

            # 翻译文本控件内文字
            if slide.has_notes_slide:
                text_frame = slide.notes_slide.notes_text_frame
                if len(text_frame.text) > 0:
                    try:
                        print("翻译前:"+text_frame.text)
                        response = self.translate_text(text=text_frame.text)
                        print("翻译后:"+response)
                        slide.notes_slide.notes_text_frame.text = response
                    except:
                        print('Invalid text in text comment ' + text_frame.text)

            # 翻译图形类控件文字,踩坑:
            # 1. 不能直接使用shape.text_frame.text,因为有的shape没有text_frame属性
            # 2. paragraph是整段的，其中run是一个样式字符块，所以可以逐个run翻译，也可以整段翻译
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    # 方式1：直接将整段话一起翻译，这样对LLM更加友好，但会损失原段落中部分文字的加粗、斜体等样式
                    try:
                        print("翻译前:"+paragraph.text)
                        response = self.translate_text(text=paragraph.text)
                        print("翻译后:"+response)
                        paragraph.text = response
                    except:
                        print('Invalid text in shapes ' + paragraph.text)

                    # 方式2：逐个run翻译，原来整句话会被拆散，LLM很可能联系不了上下文，如果对格式有要求，采用这种方式，则需要悉心调教Promot
                    # for index, paragraph_run in enumerate(paragraph.runs):
                    #     try:
                    #         print("翻译前:" + paragraph_run.text)
                    #         response = self.translate_text(text=paragraph_run.text)
                    #         print("翻译后:" + response)
                    #         paragraph.runs[index].text = response
                    #         paragraph.runs[index].font.language_id = LANGUAGE_CODE_TO_LANGUAGE_ID[self.target_language]
                    #
                    #     except:
                    #         print('Invalid text in shapes ' + paragraph_run.text)

    def translate_presentation_and_save_new(self, input_file_path):
        try:
            output_file_path = input_file_path.replace(
                '.ppt', '-[{language_code}].ppt'.format(language_code=self.target_language))
            presentation = Presentation(input_file_path)
            self.translate_presentation(presentation)
            presentation.save(output_file_path)
            print(f"Presentation saved to {output_file_path}")
            return '翻译完成,新文件路径:'+output_file_path
        except Exception as e:
            print(e)
            return e


if __name__ == '__main__':
    translator = TranslatorByXfChat(app_id='d83fbada',
                                    api_key='27438b758e01141255604a15b41b6ac3',
                                    api_secret='YWFhZjFlM2JjOTNlZjEwZmFiMTU3OWM5',
                                    version=1.1,
                                    target_language='日语'
                                    )
    print(translator.translate_text("VDM系统是一个试验管理系统"))
